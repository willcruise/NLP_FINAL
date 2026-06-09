#!/usr/bin/env python3
"""Build a PPTX with run_three_experiments.sh setting tables (native PPT tables)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HEADER_BG = RGBColor(0x2F, 0x54, 0x96)
HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ALT = RGBColor(0xEE, 0xF2, 0xFA)
ROW_BG = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)


def add_table(slide, rows, top=Inches(1.2), col_widths=None, font_size=14):
    n_rows = len(rows)
    n_cols = len(rows[0])
    left = Inches(0.5)
    width = Inches(12.3)
    height = Inches(0.4 * n_rows)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(font_size)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = HEADER_BG
                para.font.bold = True
                para.font.color.rgb = HEADER_FG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_ALT if r % 2 == 0 else ROW_BG
                para.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
            if c != 0:
                para.alignment = PP_ALIGN.CENTER
    return table


# Slide 1: common settings
s1 = prs.slides.add_slide(blank)
add_title(s1, "run_three_experiments.sh \u2013 \uacf5\ud1b5 \uc138\ud305 (exp1/2/3)")
common = [
    ["\ud56d\ubaa9", "\uac12"],
    ["Batch size", "8"],
    ["Learning rate", "5e-6"],
    ["Optimizer", "AdamW"],
    ["\ucd08\uae30\ud654", "\uc21c\uc218 pretrained GPT-2 (124M)"],
    ["--fresh", "True (\uae30\uc874 \uccb4\ud06c\ud3ec\uc778\ud2b8 \uc0ad\uc81c \ud6c4 \uc2dc\uc791)"],
    ["mask_prompt", "True, mask_target='reasoning'"],
    ["Eval \uc8fc\uae30", "2 epoch\ub9c8\ub2e4"],
    ["Early stopping", "patience 5"],
    ["Dev set", "data/multiarith_dev.jsonl (n=90)"],
    ["\uc0dd\uc131 \ud30c\ub77c\ubbf8\ud130", "temperature 0.7, top_p 0.9, max_new_tokens 256, sampling"],
    ["Seed", "11711"],
]
add_table(s1, common, col_widths=[3.5, 8.8])

# Slide 2: per-experiment data & epochs
s2 = prs.slides.add_slide(blank)
add_title(s2, "\uc2e4\ud5d8\ubcc4 \ucc28\uc774 (epoch & \ub370\uc774\ud130)")
per_exp = [
    ["\uc2e4\ud5d8", "\uc124\uc815 epoch", "\ud559\uc2b5 \ub370\uc774\ud130 \uad6c\uc131", "\ucd1d \uc608\uc2dc \uc218"],
    ["exp1", "30", "GSM8K 3,000", "3,000"],
    ["exp2", "30", "GSM8K 3,000 + MultiArith 543", "3,543"],
    ["exp3", "32", "GSM8K 3,000 + MultiArith 543 + Entity(stage2) 3,000", "6,543"],
]
add_table(s2, per_exp, col_widths=[1.4, 1.8, 7.3, 1.8], font_size=14)

# Slide 3: actual early-stopping results
s3 = prs.slides.add_slide(blank)
add_title(s3, "\uc2e4\uc81c \ud559\uc2b5 \uacb0\uacfc (early stopping)")
results = [
    ["\uc2e4\ud5d8", "best epoch", "\uc2e4\uc81c \uc885\ub8cc epoch", "best MultiArith acc"],
    ["exp1", "4", "14", "0.033"],
    ["exp2", "28", "29", "0.133"],
    ["exp3", "30", "31", "0.289"],
]
add_table(s3, results, col_widths=[2.5, 2.8, 3.2, 3.8])

out = "outputs/run_three_experiments_settings.pptx"
prs.save(out)
print("saved", out)
